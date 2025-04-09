import streamlit as st
import pytesseract
from PIL import Image
import io
from transformers import AutoTokenizer, AutoModelForCausalLM
import torch
import docx
from pptx import Presentation
import fitz  # PyMuPDF
import re

st.title("LLaMA 4: Decision Bias MCQ Tutor – Documents, Screenshots & Auto-Detection")

st.markdown("""
Upload or paste a screenshot, or submit a document (TXT, PDF, PPTX, Word, PNG, JPG).
The tutor will auto-detect MCQs and highlight the correct answer.
""")

def extract_text(file):
    ext = file.name.lower()
    if ext.endswith(".txt"):
        return file.read().decode("utf-8")
    elif ext.endswith(".docx"):
        doc = docx.Document(file)
        return "\n".join([p.text for p in doc.paragraphs])
    elif ext.endswith(".pptx"):
        prs = Presentation(file)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif ext.endswith(".pdf"):
        text = ""
        pdf = fitz.open(stream=file.read(), filetype="pdf")
        for page in pdf:
            text += page.get_text()
        return text
    elif ext.endswith((".png", ".jpg", ".jpeg")):
        image = Image.open(file)
        return pytesseract.image_to_string(image)
    else:
        return "Unsupported format"

def highlight_correct_option(text):
    # Highlight the correct option if format is 'Correct Option: X'
    match = re.search(r"Correct Option: ([A-D])", text)
    if not match:
        return text
    correct = match.group(1)
    for option in ["A", "B", "C", "D"]:
        if f"{option}." in text:
            line = re.search(rf"({option}\.\s.*)", text)
            if line:
                colored = f"<span style='background-color: #DFF2BF'><strong>{line.group(1)}</strong></span>"
                text = text.replace(line.group(1), colored)
    return text

screenshot_input = st.file_uploader("Paste or upload a screenshot or document", type=["txt", "pdf", "pptx", "docx", "png", "jpg", "jpeg"])
text_input = ""

if screenshot_input:
    text_input = extract_text(screenshot_input)

mcq_question = st.text_area("Paste your MCQ or extracted text here:", value=text_input)

if st.button("Get Answer") and mcq_question:
    tokenizer = AutoTokenizer.from_pretrained("./finetuned_model")
    model = AutoModelForCausalLM.from_pretrained("./finetuned_model")

    prompt = f"""\nYou are an expert in Decision and Behavioural Analytics. Answer the following multiple-choice question with detailed reasoning:\n\n{mcq_question}\n\nGive your answer as: Explanation followed by 'Correct Option: X'\n"""

    inputs = tokenizer(prompt, return_tensors="pt").to(model.device)
    output = model.generate(**inputs, max_new_tokens=500)
    result = tokenizer.decode(output[0], skip_special_tokens=True)
    st.markdown(highlight_correct_option(result), unsafe_allow_html=True)
